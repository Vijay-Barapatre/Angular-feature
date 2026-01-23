import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt

TRANSCRIPT_FILE = 'voiceover-script.md'
OUTPUT_PPT = 'input_output_tutorial.pptx'

def parse_transcript(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    segments = []
    # Regex to split by frames
    blocks = re.split(r'## (?:Frame|Segment)', content)
    
    for block in blocks[1:]: # Skip header
        # Extract Image path
        img_match = re.search(r'\*\*Image(?: to use)?:\*\*\s*`?([^`\n\r]+)`?', block)
        
        # Extract Transcript/Voiceover
        transcript_match = re.search(r'\*\*(?:Voiceover|Transcript):\*\*\s*(?:")?([\s\S]+?)(?:")?(?:\n---|#|$)', block)
        
        if img_match:
            img_path = img_match.group(1).strip()
            
            # Handle v2_final path correction
            if not os.path.exists(img_path) and os.path.exists(os.path.join("v2_final", img_path)):
                img_path = os.path.join("v2_final", img_path)
            elif not os.path.exists(img_path) and "v2_final/" in img_path:
                 # Check if we are already inside a related folder or need base path
                 pass

            # Clean text
            text = ""
            if transcript_match:
                 raw_text = transcript_match.group(1).strip()
                 if raw_text.startswith('"') and raw_text.endswith('"'):
                    raw_text = raw_text[1:-1]
                 text = raw_text
            
            segments.append({
                'image': img_path,
                'text': text
            })
            
    return segments

def create_ppt():
    print("Creating PowerPoint presentation...")
    segments = parse_transcript(TRANSCRIPT_FILE)
    
    if not segments:
        print("No segments found!")
        return

    # 16:9 Defaults
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6] # 6 is usually blank

    for i, seg in enumerate(segments):
        img_path = seg['image']
        text = seg['text']
        
        if not os.path.exists(img_path):
            print(f"Warning: Image not found {img_path}, skipping.")
            continue
            
        print(f"Adding Slide {i+1}: {img_path}")
        
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add Image covering the whole slide
        # left, top, width, height
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
        
        # Add Speaker Notes
        if text:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = text

    prs.save(OUTPUT_PPT)
    print(f"Successfully saved presentation to {OUTPUT_PPT}")

if __name__ == "__main__":
    create_ppt()
