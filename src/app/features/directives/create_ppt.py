"""
Angular Directives PowerPoint Generator
Creates a professional presentation about Angular Directives
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background using default solid fill
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.line.fill.background()
    
    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    accent.fill.solid()
    accent.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, bullet_points):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    
    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(24)
        p.space_after = Pt(12)
    
    return slide

def add_code_slide(title, code, description=""):
    """Add a slide with code snippet"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    y_offset = 1.4
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(12.333), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(18)
        p.font.italic = True
        y_offset = 1.9
    
    # Code box
    code_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y_offset), Inches(12.5), Inches(5.2))
    code_bg.fill.solid()
    code_bg.line.fill.background()
    
    code_box = slide.shapes.add_textbox(Inches(0.6), Inches(y_offset + 0.2), Inches(12.1), Inches(4.8))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(13)
    p.font.name = "Consolas"
    
    return slide

def add_two_column_slide(title, left_content, right_content, left_title="", right_title=""):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Left header
    if left_title:
        left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.9), Inches(0.5))
        tf = left_header.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(22)
        p.font.bold = True
    
    # Left content
    y_start = 1.9 if left_title else 1.4
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_start), Inches(5.9), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    # Right header
    if right_title:
        right_header = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.9), Inches(0.5))
        tf = right_header.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(22)
        p.font.bold = True
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(y_start), Inches(5.9), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_after = Pt(8)
    
    return slide

def add_table_slide(title, headers, rows):
    """Add a table slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_width = Inches(12)
    table_height = Inches(0.55 * num_rows)
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.65), Inches(1.5), table_width, table_height).table
    
    # Header
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
    
    # Data
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
    
    return slide

# ============================================================================
# CREATE SLIDES
# ============================================================================

# Slide 1: Title
add_title_slide("Angular Directives", "A Complete Guide to Building Powerful, Reusable Behaviors")

# Slide 2: Agenda
add_content_slide("Agenda", [
    "What Are Directives?",
    "Three Types of Directives",
    "Attribute Directives - Changing Appearance & Behavior",
    "Structural Directives - Manipulating the DOM",
    "Key APIs: ElementRef, Renderer2, TemplateRef, ViewContainerRef",
    "@HostListener & @HostBinding",
    "Real-World Production Patterns",
    "Best Practices & Performance Tips"
])

# Slide 3: What Are Directives
add_content_slide("What Are Directives?", [
    "Directives are classes that add behavior to elements in Angular",
    'Think of them as "HTML attribute superpowers"',
    "They extend what HTML elements can do!",
    "Built-in examples: *ngIf, *ngFor, [ngClass], [ngStyle]",
    "Custom directives let you create reusable behaviors"
])

# Slide 4: Three Types
add_table_slide("Three Types of Directives", ["Type", "Description", "Example"], [
    ["Components", "Directives with a template", "@Component"],
    ["Attribute", "Modify appearance/behavior", "[ngClass], [ngStyle]"],
    ["Structural", "Add/remove DOM elements", "*ngIf, *ngFor"]
])

# Slide 5: Simple Attribute Directive
add_code_slide("Simple Attribute Directive",
'''@Directive({
    selector: '[appHighlight]',
    standalone: true
})
export class HighlightDirective implements OnInit {
    private el = inject(ElementRef);
    private renderer = inject(Renderer2);

    ngOnInit(): void {
        this.renderer.setStyle(
            this.el.nativeElement, 
            'backgroundColor', 
            '#ffeb3b'
        );
    }
}

// Usage: <span appHighlight>Highlighted!</span>''',
"Creating a simple directive that highlights elements")

# Slide 6: Directive with @Input
add_code_slide("Configurable Directive with @Input",
'''@Directive({ selector: '[appHighlight]', standalone: true })
export class HighlightDirective {
    private el = inject(ElementRef);
    private renderer = inject(Renderer2);

    @Input() set appHighlight(color: string) {
        this.renderer.setStyle(
            this.el.nativeElement,
            'backgroundColor',
            color || '#ffeb3b'
        );
    }
}

// <p [appHighlight]="'yellow'">Yellow</p>
// <p [appHighlight]="'lightblue'">Blue</p>''',
"Making directives flexible with @Input")

# Slide 7: Key APIs
add_two_column_slide("Key APIs: ElementRef & Renderer2",
    ["Provides direct access to host element", "el.nativeElement gives DOM element", "Direct access - use carefully!", "May not work in SSR"],
    ["Platform-agnostic DOM operations", "setStyle(), addClass(), removeClass()", "setAttribute(), listen()", "Safe for SSR & Web Workers"],
    "ElementRef - Direct DOM Access", "Renderer2 - Safe DOM Manipulation")

# Slide 8: @HostListener
add_code_slide("@HostListener - Responding to Events",
'''@Directive({ selector: '[appHoverEffect]', standalone: true })
export class HoverEffectDirective {
    private el = inject(ElementRef);
    private renderer = inject(Renderer2);
    @Input() hoverBg = '#667eea';

    @HostListener('mouseenter')
    onMouseEnter(): void {
        this.renderer.setStyle(this.el.nativeElement, 'backgroundColor', this.hoverBg);
        this.renderer.setStyle(this.el.nativeElement, 'transform', 'scale(1.05)');
    }

    @HostListener('mouseleave')
    onMouseLeave(): void {
        this.renderer.removeStyle(this.el.nativeElement, 'backgroundColor');
    }
}''',
"Listen to host element events declaratively")

# Slide 9: @HostBinding
add_code_slide("@HostBinding - Binding Properties",
'''@Directive({ selector: '[appActiveToggle]', standalone: true })
export class ActiveToggleDirective {
    private isActive = false;

    @HostBinding('class.active')
    get active(): boolean { return this.isActive; }

    @HostBinding('style.backgroundColor')
    get bgColor(): string { 
        return this.isActive ? '#4ade80' : '#f87171'; 
    }

    @HostListener('click')
    toggle(): void { this.isActive = !this.isActive; }
}''',
"Bind host element properties and attributes")

# Slide 10: Structural Directives
add_content_slide("Structural Directives", [
    "Change the DOM structure by adding/removing elements",
    "Prefixed with asterisk (*) - syntactic sugar",
    "*ngIf='condition' expands to <ng-template [ngIf]='condition'>",
    "Key APIs: TemplateRef (blueprint) + ViewContainerRef (slot)",
    "createEmbeddedView() to stamp, clear() to remove"
])

# Slide 11: TemplateRef & ViewContainerRef
add_two_column_slide("TemplateRef & ViewContainerRef",
    ["The 'Rubber Stamp'", "Holds the template blueprint", "Contains HTML inside *directive", "Doesn't render by itself"],
    ["The 'Slot on Page'", "Location to render template", "createEmbeddedView(templateRef)", "clear() removes all views"],
    "TemplateRef<T>", "ViewContainerRef")

# Slide 12: Custom *appIf
add_code_slide("Custom Structural Directive: *appIf",
'''@Directive({ selector: '[appIf]', standalone: true })
export class AppIfDirective implements OnChanges {
    private templateRef = inject(TemplateRef<any>);
    private viewContainer = inject(ViewContainerRef);
    private hasView = false;
    @Input() appIf = false;

    ngOnChanges(): void {
        if (this.appIf && !this.hasView) {
            this.viewContainer.createEmbeddedView(this.templateRef);
            this.hasView = true;
        } else if (!this.appIf && this.hasView) {
            this.viewContainer.clear();
            this.hasView = false;
        }
    }
}''',
"Building your own *ngIf equivalent")

# Slide 13: Permission Directive
add_code_slide("Real-World: Permission-Based Visibility",
'''@Directive({ selector: '[appPermission]', standalone: true })
export class PermissionDirective implements OnInit {
    private templateRef = inject(TemplateRef<any>);
    private viewContainer = inject(ViewContainerRef);
    @Input() appPermission: string[] = [];
    private currentUserRoles = ['user', 'editor'];

    ngOnInit(): void {
        const hasPermission = this.appPermission.some(role =>
            this.currentUserRoles.includes(role)
        );
        if (hasPermission) {
            this.viewContainer.createEmbeddedView(this.templateRef);
        }
    }
}
// <button *appPermission="['admin']">Delete</button>''',
"Role-based access control in templates")

# Slide 14: Lazy Load
add_code_slide("Real-World: Lazy Load Images",
'''@Directive({ selector: '[appLazyLoad]', standalone: true })
export class LazyLoadDirective implements AfterViewInit {
    private el = inject(ElementRef);
    private observer: IntersectionObserver | null = null;
    @Input() appLazyLoad = '';
    @Output() loaded = new EventEmitter<void>();

    ngAfterViewInit(): void {
        this.observer = new IntersectionObserver(entries => {
            if (entries[0].isIntersecting) {
                this.el.nativeElement.src = this.appLazyLoad;
                this.loaded.emit();
                this.observer?.disconnect();
            }
        }, { threshold: 0.1 });
        this.observer.observe(this.el.nativeElement);
    }
}''',
"Load images only when they enter viewport")

# Slide 15: Directive Catalog
add_table_slide("Directive Catalog Summary", ["Directive", "Type", "Use Case"], [
    ["[appHighlight]", "Attribute", "Apply background color"],
    ["[appHoverEffect]", "Attribute", "Mouse hover effects"],
    ["[appCopyToClipboard]", "Attribute", "Copy text on click"],
    ["[appDebounceClick]", "Attribute", "Prevent double clicks"],
    ["[appLazyLoad]", "Attribute", "Lazy load images"],
    ["*appIf", "Structural", "Conditional rendering"],
    ["*appPermission", "Structural", "Role-based visibility"]
])

# Slide 16: Best Practices
add_two_column_slide("Best Practices",
    ["Use Renderer2 for DOM (SSR-safe)", "Use inject() for DI", "Clean up in ngOnDestroy", "Use @Input setters for reactivity", "Keep directives focused", "Use standalone: true"],
    ["Direct nativeElement for styling", "Forget to unsubscribe", "Create multi-purpose directives", "Use directives when component fits", "Ignore memory leaks", "Skip error handling"],
    "DO", "DON'T")

# Slide 17: Lifecycle Hooks
add_table_slide("Lifecycle Hook Selection", ["Hook", "Use When"], [
    ["constructor()", "Dependencies ready, no inputs needed"],
    ["ngOnInit()", "Inputs ready, set up logic/listeners"],
    ["ngOnChanges()", "React to input changes"],
    ["ngAfterViewInit()", "DOM fully rendered, focus/scroll"],
    ["ngOnDestroy()", "Cleanup observers/listeners"]
])

# Slide 18: Key Takeaways
add_content_slide("Key Takeaways", [
    "Directives extend HTML with custom behaviors",
    "Attribute directives modify appearance/behavior",
    "Structural directives change DOM structure",
    "Use Renderer2 for safe DOM manipulation",
    "TemplateRef + ViewContainerRef = Structural magic",
    "@HostListener for events, @HostBinding for properties",
    "Always clean up in ngOnDestroy!"
])

# Slide 19: Thank You
add_title_slide("Thank You!", "Directives are the secret sauce that makes Angular templates powerful.")

# Save
output_path = os.path.join(os.path.dirname(__file__), "Angular_Directives_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
